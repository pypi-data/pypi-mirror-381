import os
import platform
import shutil
import subprocess
from abc import ABC, abstractmethod
from pathlib import Path
from typing import List, Union

from pptx import Presentation

from pptx_to_video.logger import logger


class PowerPointEngine(ABC):
    """Interface for PowerPoint engines"""

    def __init__(self, pptx_path):
        self.pptx_path = pptx_path

    @abstractmethod
    def export_slides_as_images(self, out_dir: Union[str, os.PathLike]) -> List[Path]:
        """Saves slides as images to out_dir

        :param out_dir: Output directory
        :type out_dir: Union[str,os.PathLike]
        :return: List of image paths
        :rtype: List[Path]
        """
        pass

    def extract_notes_from_pptx(self) -> List[dict]:
        """Extract notes from pptx slides

        :return: List of dicts with slide index and notes
        :rtype: List[dict]
        """
        prs = Presentation(self.pptx_path)
        slides_data = []
        for i, slide in enumerate(prs.slides):
            logger.info(f"Extracting notes from Slide {i + 1}/{len(prs.slides)}")
            notes_text = ""
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = "\n".join(
                    p.text for p in slide.notes_slide.notes_text_frame.paragraphs
                ).strip()
            slides_data.append({"index": i, "notes": notes_text or ""})

        logger.info("Finished extracting notes from PPTX")
        return slides_data


class WindowsPowerPointEngine(PowerPointEngine):
    """PowerPoint engine using local PowerPoint application on Windows"""

    def __init__(self, pptx_path):
        super().__init__(pptx_path)

        if platform.system() != "Windows":
            raise EnvironmentError(
                "WindowsPowerPointEngine can only be used on Windows systems."
            )

        if not pptx_path.lower().endswith((".pptx", ".ppt")):
            raise ValueError(
                "The provided file is not a valid PowerPoint file (.pptx or .ppt)"
            )

        if not os.path.isfile(pptx_path):
            raise FileNotFoundError(f"The file {pptx_path} does not exist.")

    def export_slides_as_images(self, out_dir: Union[str, os.PathLike]) -> List[Path]:
        """Saves slides as images to out_dir

        :param out_dir: Output directory
        :type out_dir: Union[str,os.PathLike]
        :return: List of image paths
        :rtype: List[Path]
        """
        import win32com.client

        out_dir = Path(out_dir)
        out_dir.mkdir(parents=True, exist_ok=True)

        logger.info(f"Exporting slides as images to {out_dir}...")

        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        powerpoint.Visible = True
        presentation = powerpoint.Presentations.Open(
            os.path.abspath(self.pptx_path), WithWindow=False
        )
        base = out_dir / "slide_image"
        presentation.SaveAs(str(base), 17)  # 17 is the enum for JPG format
        presentation.Close()
        powerpoint.Quit()

        images = sorted(out_dir.glob("slide_image*.JPG"))
        logger.info(f"Exported {len(images)} slides as images.")
        return images


class LibreOfficePowerPointEngine(PowerPointEngine):
    """PowerPoint engine using LibreOffice on Linux and macOS"""

    def __init__(self, pptx_path):
        super().__init__(pptx_path)

        if platform.system() not in ["Linux", "Darwin"]:
            raise EnvironmentError(
                "LibreOfficePowerPointEngine can only be used on Linux or macOS systems."
            )

        if not pptx_path.lower().endswith((".pptx", ".ppt")):
            raise ValueError(
                "The provided file is not a valid PowerPoint file (.pptx or .ppt)"
            )

        if not os.path.isfile(pptx_path):
            raise FileNotFoundError(f"The file {pptx_path} does not exist.")

        # Check if libreoffice command is available
        if not shutil.which("libreoffice"):
            raise EnvironmentError(
                "LibreOffice is not installed or not found in system PATH."
            )

    def export_slides_as_images(self, out_dir: Union[str, os.PathLike]) -> List[Path]:
        """Saves slides as images to out_dir

        :param out_dir: Output directory
        :type out_dir: Union[str,os.PathLike]
        :return: List of image paths
        :rtype: List[Path]
        """
        out_dir = Path(out_dir)
        out_dir.mkdir(parents=True, exist_ok=True)

        logger.info(f"Exporting slides as images to {out_dir} using LibreOffice...")
        try:
            subprocess.run(
                [
                    "libreoffice",
                    "--headless",
                    "--convert-to",
                    "jpg",
                    "--outdir",
                    str(out_dir),
                    str(self.pptx_path),
                ],
                check=True,
            )
        except subprocess.CalledProcessError as e:
            raise RuntimeError(
                "Failed to convert PPTX to images using LibreOffice."
            ) from e

        images = sorted(out_dir.glob("*.jpg"))
        logger.info(f"Exported {len(images)} slides as images.")
        return images
