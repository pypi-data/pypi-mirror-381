from __future__ import annotations
from typing import Iterable
from ..types import UnifiedDocument, Metadata, Section, Chunk
from ..core import register_parser, ParserProtocol

class PPTXParser(ParserProtocol):
    name = "pptx"
    content_types: Iterable[str] = ("application/vnd.openxmlformats-officedocument.presentationml.presentation",)
    extensions: Iterable[str] = (".pptx",)

    def can_parse(self, meta: Metadata) -> bool:
        return (meta.content_type in self.content_types or 
                (meta.path or "").endswith(".pptx"))

    def parse(self, target, meta: Metadata, recursive: bool = False, **kwargs) -> UnifiedDocument:
        doc = UnifiedDocument(meta=meta, sections=[])
        
        try:
            from pptx import Presentation
            
            # Load presentation
            prs = Presentation(str(target))
            
            # Extract presentation properties
            if prs.core_properties.title:
                doc.meta.title = prs.core_properties.title
            if prs.core_properties.author:
                doc.meta.extra["author"] = prs.core_properties.author
            if prs.core_properties.subject:
                doc.meta.extra["subject"] = prs.core_properties.subject
            if prs.core_properties.created:
                doc.meta.extra["created"] = prs.core_properties.created.isoformat()
            if prs.core_properties.modified:
                doc.meta.extra["modified"] = prs.core_properties.modified.isoformat()
            
            order = 0
            
            # Process each slide
            for i, slide in enumerate(prs.slides, 1):
                slide_texts = []
                slide_title = None
                
                # Extract text from all shapes on the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text = shape.text.strip()
                        
                        # Check if this might be a title (usually the first text box or larger text)
                        if not slide_title and (shape.shape_type == 1 or "title" in str(shape).lower()):
                            slide_title = text
                        else:
                            slide_texts.append(text)
                
                # Create section for this slide
                if slide_texts or slide_title:
                    all_text = []
                    if slide_title:
                        all_text.append(f"Title: {slide_title}")
                    all_text.extend(slide_texts)
                    
                    slide_content = "\n".join(all_text)
                    doc.sections.append(Section(
                        heading=f"Slide {i}" + (f": {slide_title}" if slide_title else ""),
                        chunks=[Chunk(text=slide_content, order=order, meta={"type": "slide"})],
                        meta={"type": "slide", "slide_number": i, "title": slide_title}
                    ))
                    order += 1
            
            # Store presentation info
            doc.meta.extra["presentation_info"] = {
                "slide_count": len(prs.slides),
                "has_images": any(shape.shape_type == 13 for slide in prs.slides for shape in slide.shapes)
            }
            
        except ImportError:
            error_text = "[panparsex:pptx] python-pptx not available. Install with: pip install python-pptx"
            doc.sections.append(Section(
                heading="Error",
                chunks=[Chunk(text=error_text, order=0, meta={"error": True})],
                meta={"error": True}
            ))
        except Exception as e:
            error_text = f"[panparsex:pptx] unable to parse: {e}"
            doc.sections.append(Section(
                heading="Error",
                chunks=[Chunk(text=error_text, order=0, meta={"error": True})],
                meta={"error": True}
            ))
        
        return doc

register_parser(PPTXParser())
